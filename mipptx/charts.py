from __future__ import annotations

from io import BytesIO
from typing import List, Literal, Optional

import openpyxl
from pptx.chart.data import BubbleChartData, CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE as XL
from pydantic import ConfigDict, Field

from .base import JsonModel
from .enums import ShapeKind
from .utils import emu_to_pt, pt_to_emu


class ChartType(str):
    # Values match python-pptx XL_CHART_TYPE member names (lowercased)
    COLUMN_CLUSTERED = "column_clustered"
    BAR_CLUSTERED = "bar_clustered"
    LINE = "line"
    LINE_MARKERS = "line_markers"
    PIE = "pie"
    DOUGHNUT = "doughnut"
    AREA = "area"
    SCATTER = "scatter"
    BUBBLE = "bubble"


def _to_chart_type(xl_type) -> str:
    mapping = {
        XL.COLUMN_CLUSTERED: ChartType.COLUMN_CLUSTERED,
        XL.BAR_CLUSTERED: ChartType.BAR_CLUSTERED,
        XL.LINE: ChartType.LINE,
        XL.LINE_MARKERS: ChartType.LINE_MARKERS,
        XL.PIE: ChartType.PIE,
        XL.DOUGHNUT: ChartType.DOUGHNUT,
        XL.AREA: ChartType.AREA,
        # XY scatter variants (name differences across versions)
        getattr(XL, "XY_SCATTER", None): ChartType.SCATTER,
        getattr(XL, "XY_SCATTER_LINES", None): ChartType.SCATTER,
        getattr(XL, "XY_SCATTER_LINES_NO_MARKERS", None): ChartType.SCATTER,
        getattr(XL, "XY_SCATTER_SMOOTH", None): ChartType.SCATTER,
        getattr(XL, "XY_SCATTER_SMOOTH_NO_MARKERS", None): ChartType.SCATTER,
        XL.BUBBLE: ChartType.BUBBLE,
        XL.BUBBLE_THREE_D_EFFECT: ChartType.BUBBLE,
    }
    mapping = {k: v for k, v in mapping.items() if k is not None}
    return mapping.get(xl_type, ChartType.COLUMN_CLUSTERED)


def _from_chart_type(chart_type: str):
    reverse = {
        ChartType.COLUMN_CLUSTERED: XL.COLUMN_CLUSTERED,
        ChartType.BAR_CLUSTERED: XL.BAR_CLUSTERED,
        ChartType.LINE: XL.LINE,
        ChartType.LINE_MARKERS: XL.LINE_MARKERS,
        ChartType.PIE: XL.PIE,
        ChartType.DOUGHNUT: XL.DOUGHNUT,
        ChartType.AREA: XL.AREA,
        ChartType.SCATTER: getattr(XL, "XY_SCATTER", None),
        ChartType.BUBBLE: XL.BUBBLE,
    }
    return reverse.get(chart_type)


class CategorySeriesModel(JsonModel):
    model_config = ConfigDict(extra="forbid")

    name: str = ""
    values: List[float] = Field(default_factory=list)
    number_format: Optional[str] = None


class CategoryChartDataModel(JsonModel):
    model_config = ConfigDict(extra="forbid")

    categories: List[str] = Field(default_factory=list)
    series: List[CategorySeriesModel] = Field(default_factory=list)

    @classmethod
    def from_chart(cls, chart) -> "CategoryChartDataModel":
        # Use first plot for categories
        cats: List[str] = []
        try:
            plot0 = chart.plots[0]
            cats = [str(c) for c in plot0.categories]
        except Exception:
            cats = []

        series: List[CategorySeriesModel] = []
        try:
            for s in chart.series:
                series.append(
                    CategorySeriesModel(name=s.name or "", values=list(s.values))
                )
        except Exception:
            pass

        return cls(categories=cats, series=series)

    def to_chart_data(self):
        cd = CategoryChartData()
        if self.categories:
            cd.categories = list(self.categories)
        for s in self.series:
            cd.add_series(s.name, values=s.values, number_format=s.number_format)
        return cd


class XyPoint(JsonModel):
    x: float | None = None
    y: float | None = None


class XySeriesModel(JsonModel):
    model_config = ConfigDict(extra="forbid")
    name: str = ""
    points: list[XyPoint] = Field(default_factory=list)
    number_format: Optional[str] = None


class XyChartDataModel(JsonModel):
    model_config = ConfigDict(extra="forbid")
    series: list[XySeriesModel] = Field(default_factory=list)

    @classmethod
    def from_chart(cls, chart) -> "XyChartDataModel":
        out_series: list[XySeriesModel] = []
        try:
            for s in chart.series:
                name = s.name or ""
                pts: list[XyPoint] = []
                try:
                    ser = getattr(s, "_ser", None) or getattr(s, "_element", None)
                    if (
                        ser is not None
                        and getattr(ser, "xVal", None) is not None
                        and getattr(ser, "yVal", None) is not None
                    ):
                        count = min(ser.xVal.ptCount_val, ser.yVal.ptCount_val)
                        for i in range(count):
                            x = ser.xVal.pt_v(i)
                            y = ser.yVal.pt_v(i)
                            pts.append(XyPoint(x=x, y=y))
                    else:
                        ys = list(getattr(s, "values", []) or [])
                        pts = [
                            XyPoint(x=float(i), y=float(v) if v is not None else None)
                            for i, v in enumerate(ys)
                        ]
                except Exception:
                    pass
                out_series.append(XySeriesModel(name=name, points=pts))
        except Exception:
            pass
        return cls(series=out_series)

    def to_chart_data(self):
        cd = XyChartData()
        for s in self.series:
            xs = cd.add_series(s.name, number_format=s.number_format)
            for p in s.points:
                xs.add_data_point(p.x, p.y)
        return cd


class BubblePoint(JsonModel):
    x: float | None = None
    y: float | None = None
    size: float | None = None


class BubbleSeriesModel(JsonModel):
    model_config = ConfigDict(extra="forbid")
    name: str = ""
    points: list[BubblePoint] = Field(default_factory=list)
    number_format: Optional[str] = None


class BubbleChartDataModel(JsonModel):
    model_config = ConfigDict(extra="forbid")
    series: list[BubbleSeriesModel] = Field(default_factory=list)

    @classmethod
    def from_chart(cls, chart) -> "BubbleChartDataModel":
        out_series: list[BubbleSeriesModel] = []
        try:
            for s in chart.series:
                name = s.name or ""
                pts: list[BubblePoint] = []
                try:
                    ser = getattr(s, "_ser", None) or getattr(s, "_element", None)
                    if (
                        ser is not None
                        and getattr(ser, "xVal", None) is not None
                        and getattr(ser, "yVal", None) is not None
                        and getattr(ser, "bubbleSize", None) is not None
                    ):
                        count = min(
                            ser.xVal.ptCount_val,
                            ser.yVal.ptCount_val,
                            ser.bubbleSize_ptCount_val,
                        )
                        for i in range(count):
                            x = ser.xVal.pt_v(i)
                            y = ser.yVal.pt_v(i)
                            try:
                                matches = ser.xpath("./c:bubbleSize//c:pt[@idx=%d]" % i)
                                size = matches[0].value if matches else None
                            except Exception:
                                size = None
                            pts.append(BubblePoint(x=x, y=y, size=size))
                    else:
                        ys = list(getattr(s, "values", []) or [])
                        pts = [
                            BubblePoint(
                                x=float(i),
                                y=float(v) if v is not None else None,
                                size=None,
                            )
                            for i, v in enumerate(ys)
                        ]
                except Exception:
                    pass
                out_series.append(BubbleSeriesModel(name=name, points=pts))
        except Exception:
            pass
        return cls(series=out_series)

    def to_chart_data(self):
        cd = BubbleChartData()
        for s in self.series:
            bs = cd.add_series(s.name, number_format=s.number_format)
            for p in s.points:
                bs.add_data_point(p.x, p.y, p.size)
        return cd


class ChartModel(JsonModel):
    model_config = ConfigDict(extra="forbid")

    type: Literal[ShapeKind.chart] = ShapeKind.chart
    name: Optional[str] = None
    left_pt: float = Field(...)
    top_pt: float = Field(...)
    width_pt: float = Field(..., ge=0)
    height_pt: float = Field(..., ge=0)
    rotation: Optional[float] = 0.0

    chart_type: str = Field(ChartType.COLUMN_CLUSTERED)
    title: Optional[str] = None
    has_legend: Optional[bool] = None
    style: Optional[int] = Field(None, ge=1, le=48)

    # Data
    category_data: Optional[CategoryChartDataModel] = None
    xy_data: Optional[XyChartDataModel] = None
    bubble_data: Optional[BubbleChartDataModel] = None

    @classmethod
    def from_pptx(cls, gframe) -> "ChartModel":
        ch = gframe.chart
        chart_type = _to_chart_type(ch.chart_type)
        title = None
        try:
            ct = getattr(ch, "chart_title", None)
            if getattr(ch, "has_title", False) and ct is not None:
                tf = getattr(ct, "text_frame", None)
                if tf is not None:
                    title = getattr(tf, "text", None)
        except Exception:
            title = None
        has_legend = None
        try:
            has_legend = bool(ch.has_legend)
        except Exception:
            pass
        style = getattr(ch, "chart_style", None)

        # data extraction per chart family
        cat_data = None
        xy_data = None
        bubble_data = None
        try:
            if chart_type == ChartType.SCATTER:
                xy_data = XyChartDataModel.from_chart(ch)
            elif chart_type == ChartType.BUBBLE:
                bubble_data = BubbleChartDataModel.from_chart(ch)
            else:
                cat_data = CategoryChartDataModel.from_chart(ch)
        except Exception:
            pass

        return cls(
            name=getattr(gframe, "name", None),
            left_pt=emu_to_pt(gframe.left) or 0.0,
            top_pt=emu_to_pt(gframe.top) or 0.0,
            width_pt=emu_to_pt(gframe.width) or 0.0,
            height_pt=emu_to_pt(gframe.height) or 0.0,
            rotation=getattr(gframe, "rotation", 0.0) or 0.0,
            chart_type=chart_type,
            title=title,
            has_legend=has_legend,
            style=style,
            category_data=cat_data,
            xy_data=xy_data,
            bubble_data=bubble_data,
        )

    def apply_to_slide(self, slide) -> None:
        # Build a new chart from whichever data block is present
        xl_type = _from_chart_type(self.chart_type)
        if xl_type is None:
            return
        # Pick matching data object
        cd = None
        if self.category_data is not None:
            cd = self.category_data.to_chart_data()
        elif self.xy_data is not None:
            cd = self.xy_data.to_chart_data()
        elif self.bubble_data is not None:
            cd = self.bubble_data.to_chart_data()
        if cd is None:
            return
        left = pt_to_emu(self.left_pt)
        top = pt_to_emu(self.top_pt)
        width = pt_to_emu(self.width_pt)
        height = pt_to_emu(self.height_pt)
        gframe = slide.shapes.add_chart(xl_type, left, top, width, height, cd)
        # gframe is a GraphicFrame; access chart
        ch = gframe.chart
        if self.name:
            gframe.name = self.name
        if self.style is not None:
            try:
                ch.chart_style = int(self.style)
            except Exception:
                pass
        if self.has_legend is not None:
            try:
                ch.has_legend = bool(self.has_legend)
            except Exception:
                pass
        if self.title is not None:
            try:
                ch.has_title = True
                ch.chart_title.text_frame.text = self.title
            except Exception:
                pass

    # Update data and selected properties on an existing chart, with formatting snapshot/restore
    def apply_to_existing_chart(self, chart) -> None:
        # snapshot series-level data label settings to preserve formatting across data replacement
        def _snapshot_dlbls(ch):
            out = []
            try:
                for s in ch.series:
                    try:
                        dl = s.data_labels
                        out.append(
                            dict(
                                number_format=getattr(dl, "number_format", None),
                                number_format_is_linked=getattr(
                                    dl, "number_format_is_linked", True
                                ),
                                position=getattr(dl, "position", None),
                                show_value=getattr(dl, "show_value", None),
                                show_series_name=getattr(dl, "show_series_name", None),
                                show_category_name=getattr(
                                    dl, "show_category_name", None
                                ),
                                show_percentage=getattr(dl, "show_percentage", None),
                            )
                        )
                    except Exception:
                        out.append({})
            except Exception:
                pass
            return out

        def _restore_dlbls(ch, snapshot):
            try:
                for s, snap in zip(ch.series, snapshot):
                    try:
                        dl = s.data_labels
                        if snap.get("number_format") is not None:
                            dl.number_format = snap["number_format"]
                        if snap.get("number_format_is_linked") is not None:
                            dl.number_format_is_linked = snap["number_format_is_linked"]
                        if snap.get("position") is not None:
                            dl.position = snap["position"]
                        for key in (
                            "show_value",
                            "show_series_name",
                            "show_category_name",
                            "show_percentage",
                        ):
                            if snap.get(key) is not None:
                                setattr(dl, key, snap[key])
                    except Exception:
                        pass
            except Exception:
                pass

        dl_snapshot = _snapshot_dlbls(chart)

        # snapshot axes and legend settings to preserve chart formatting
        def _snap_axis(ax):
            if ax is None:
                return None
            d = {}
            try:
                d.update(
                    has_title=getattr(ax, "has_title", None),
                    title_text=(
                        ax.axis_title.text_frame.text
                        if getattr(ax, "has_title", False)
                        else None
                    ),
                    major_tick_mark=getattr(ax, "major_tick_mark", None),
                    minor_tick_mark=getattr(ax, "minor_tick_mark", None),
                    minimum_scale=getattr(ax, "minimum_scale", None),
                    maximum_scale=getattr(ax, "maximum_scale", None),
                    major_unit=getattr(ax, "major_unit", None),
                    minor_unit=getattr(ax, "minor_unit", None),
                    reverse_order=getattr(ax, "reverse_order", None),
                    tick_label_position=getattr(ax, "tick_label_position", None),
                    has_major_gridlines=getattr(ax, "has_major_gridlines", None),
                    has_minor_gridlines=getattr(ax, "has_minor_gridlines", None),
                    tick_num_fmt=(
                        getattr(getattr(ax, "tick_labels", None), "number_format", None)
                        if hasattr(ax, "tick_labels")
                        else None
                    ),
                    tick_num_fmt_linked=(
                        getattr(
                            getattr(ax, "tick_labels", None),
                            "number_format_is_linked",
                            None,
                        )
                        if hasattr(ax, "tick_labels")
                        else None
                    ),
                    visible=getattr(ax, "visible", None),
                )
            except Exception:
                pass
            return d

        def _restore_axis(ax, d):
            if ax is None or not d:
                return
            try:
                if d.get("has_title") is not None:
                    ax.has_title = d["has_title"]
                    if d["has_title"] and d.get("title_text") is not None:
                        try:
                            ax.axis_title.text_frame.text = d["title_text"]
                        except Exception:
                            pass
                for k in (
                    "major_tick_mark",
                    "minor_tick_mark",
                    "minimum_scale",
                    "maximum_scale",
                    "major_unit",
                    "minor_unit",
                    "reverse_order",
                    "tick_label_position",
                    "has_major_gridlines",
                    "has_minor_gridlines",
                    "visible",
                ):
                    if d.get(k) is not None:
                        try:
                            setattr(ax, k, d[k])
                        except Exception:
                            pass
                tl = getattr(ax, "tick_labels", None)
                if tl is not None:
                    if d.get("tick_num_fmt") is not None:
                        tl.number_format = d["tick_num_fmt"]
                    if d.get("tick_num_fmt_linked") is not None:
                        tl.number_format_is_linked = d["tick_num_fmt_linked"]
            except Exception:
                pass

        def _snapshot_axes(ch):
            snap = {}
            for name in ("category_axis", "value_axis"):
                try:
                    ax = getattr(ch, name)
                except Exception:
                    ax = None
                snap[name] = _snap_axis(ax)
            return snap

        def _restore_axes(ch, snap):
            for name in ("category_axis", "value_axis"):
                try:
                    ax = getattr(ch, name)
                except Exception:
                    ax = None
                _restore_axis(ax, snap.get(name))

        # legend snapshot
        def _snapshot_legend(ch):
            try:
                if not getattr(ch, "has_legend", False):
                    return dict(has_legend=False)
                lg = ch.legend
                return dict(
                    has_legend=True,
                    position=getattr(lg, "position", None),
                    include_in_layout=getattr(lg, "include_in_layout", None),
                )
            except Exception:
                return {}

        def _restore_legend(ch, d):
            try:
                if d.get("has_legend") is not None:
                    ch.has_legend = d["has_legend"]
                if d.get("has_legend"):
                    lg = ch.legend
                    if d.get("position") is not None:
                        lg.position = d["position"]
                    if d.get("include_in_layout") is not None:
                        lg.include_in_layout = d["include_in_layout"]
            except Exception:
                pass

        axes_snapshot = _snapshot_axes(chart)
        legend_snapshot = _snapshot_legend(chart)

        # replace data based on provided data block
        cd = None
        if self.category_data is not None:
            cd = self.category_data.to_chart_data()
        elif self.xy_data is not None:
            cd = self.xy_data.to_chart_data()
        elif self.bubble_data is not None:
            cd = self.bubble_data.to_chart_data()
        if cd is not None:
            try:
                chart.replace_data(cd)
            except Exception:
                pass
            else:
                _restore_dlbls(chart, dl_snapshot)
                _restore_axes(chart, axes_snapshot)
                _restore_legend(chart, legend_snapshot)
        # optional properties
        if self.style is not None:
            try:
                chart.chart_style = int(self.style)
            except Exception:
                pass
        if self.has_legend is not None:
            try:
                chart.has_legend = bool(self.has_legend)
            except Exception:
                pass
        if self.title is not None:
            try:
                chart.has_title = True
                chart.chart_title.text_frame.text = self.title
            except Exception:
                pass

    # Workbook-only updates: edit the embedded workbook cells referenced by the chart
    # without touching chart XML. Returns True if updated, False on fallback.
    def update_workbook_only(self, chart) -> bool:
        # helper: parse a formula like Sheet1!$B$2:$B$5 into (sheet, cells[]) in row-major order
        def col_to_idx(col: str) -> int:
            n = 0
            for ch in col:
                if ch == "$":
                    continue
                n = n * 26 + (ord(ch.upper()) - ord("A") + 1)
            return n

        def idx_to_col(n: int) -> str:
            s = ""
            while n > 0:
                n, r = divmod(n - 1, 26)
                s = chr(65 + r) + s
            return s

        def parse_formula(f: str):
            if not f:
                return None
            txt = f.strip()
            if "!" not in txt:
                return None
            sheet, rng = txt.split("!", 1)
            if sheet.startswith("'") and sheet.endswith("'"):
                sheet = sheet[1:-1]
            parts = rng.split(":")

            def parse_cell(c: str):
                c = c.replace("$", "")
                i = 0
                while i < len(c) and c[i].isalpha():
                    i += 1
                col = c[:i]
                row = int(c[i:]) if i < len(c) else 0
                return col, row

            if len(parts) == 1:
                col, row = parse_cell(parts[0])
                return sheet, [f"{col}{row}"]
            (c1, r1) = parse_cell(parts[0])
            (c2, r2) = parse_cell(parts[1])
            c1i, c2i = col_to_idx(c1), col_to_idx(c2)
            cells = []
            for r in range(min(r1, r2), max(r1, r2) + 1):
                for ci in range(min(c1i, c2i), max(c1i, c2i) + 1):
                    cells.append(f"{idx_to_col(ci)}{r}")
            return sheet, cells

        # load workbook blob
        try:
            xlsx_part = chart.part.chart_workbook.xlsx_part
            if xlsx_part is None:
                return False
            wb = openpyxl.load_workbook(BytesIO(xlsx_part.blob))
        except Exception:
            return False

        ser_list = list(chart.series)

        # category, scatter, or bubble
        def get_formula(ser, path: str):
            try:
                matches = ser._ser.xpath(path)  # type: ignore[attr-defined]
                if not matches:
                    return None
                return matches[0].text
            except Exception:
                return None

        # write helpers
        def write_range(f: str, values: list):
            pf = parse_formula(f)
            if pf is None:
                return
            sheet, cells = pf
            if sheet not in wb.sheetnames:
                return
            ws = wb[sheet]
            for i, cell in enumerate(cells):
                if i >= len(values):
                    break
                ws[cell].value = values[i]

        try:
            if self.category_data is not None:
                # categories: from first series cat ref
                if ser_list:
                    f_cat = get_formula(
                        ser_list[0], "./c:cat/c:strRef/c:f"
                    ) or get_formula(ser_list[0], "./c:cat/c:numRef/c:f")
                    if f_cat and self.category_data.categories:
                        write_range(f_cat, list(self.category_data.categories))
                # series values
                for idx, s in enumerate(ser_list):
                    f_val = get_formula(s, "./c:val/c:numRef/c:f")
                    if not f_val:
                        continue
                    vals = []
                    try:
                        vals = list(self.category_data.series[idx].values)
                    except Exception:
                        vals = []
                    write_range(f_val, vals)
            elif self.xy_data is not None:
                for idx, s in enumerate(ser_list):
                    f_x = get_formula(s, "./c:xVal/c:numRef/c:f")
                    f_y = get_formula(s, "./c:yVal/c:numRef/c:f")
                    pts = []
                    try:
                        pts = list(self.xy_data.series[idx].points)
                    except Exception:
                        pts = []
                    if f_x:
                        write_range(f_x, [p.x for p in pts])
                    if f_y:
                        write_range(f_y, [p.y for p in pts])
            elif self.bubble_data is not None:
                for idx, s in enumerate(ser_list):
                    f_x = get_formula(s, "./c:xVal/c:numRef/c:f")
                    f_y = get_formula(s, "./c:yVal/c:numRef/c:f")
                    f_b = get_formula(s, "./c:bubbleSize/c:numRef/c:f")
                    pts = []
                    try:
                        pts = list(self.bubble_data.series[idx].points)
                    except Exception:
                        pts = []
                    if f_x:
                        write_range(f_x, [p.x for p in pts])
                    if f_y:
                        write_range(f_y, [p.y for p in pts])
                    if f_b:
                        write_range(f_b, [p.size for p in pts])
            else:
                return False
        except Exception:
            return False

        # save back to blob
        try:
            bio = BytesIO()
            wb.save(bio)
            chart.part.chart_workbook.update_from_xlsx_blob(bio.getvalue())
            return True
        except Exception:
            return False

    # Update data and selected properties on an existing chart (no geometry changes)
    def apply_to_existing_chart(self, chart) -> None:
        # snapshot series-level data label settings to preserve formatting across data replacement
        def _snapshot_dlbls(ch):
            out = []
            try:
                for s in ch.series:
                    try:
                        dl = s.data_labels
                        out.append(
                            dict(
                                number_format=getattr(dl, "number_format", None),
                                number_format_is_linked=getattr(
                                    dl, "number_format_is_linked", True
                                ),
                                position=getattr(dl, "position", None),
                                show_value=getattr(dl, "show_value", None),
                                show_series_name=getattr(dl, "show_series_name", None),
                                show_category_name=getattr(
                                    dl, "show_category_name", None
                                ),
                                show_percentage=getattr(dl, "show_percentage", None),
                            )
                        )
                    except Exception:
                        out.append({})
            except Exception:
                pass
            return out

        def _restore_dlbls(ch, snapshot):
            try:
                for s, snap in zip(ch.series, snapshot):
                    try:
                        dl = s.data_labels
                        if snap.get("number_format") is not None:
                            dl.number_format = snap["number_format"]
                        if snap.get("number_format_is_linked") is not None:
                            dl.number_format_is_linked = snap["number_format_is_linked"]
                        if snap.get("position") is not None:
                            dl.position = snap["position"]
                        for key in (
                            "show_value",
                            "show_series_name",
                            "show_category_name",
                            "show_percentage",
                        ):
                            if snap.get(key) is not None:
                                setattr(dl, key, snap[key])
                    except Exception:
                        pass
            except Exception:
                pass

        dl_snapshot = _snapshot_dlbls(chart)

        # snapshot axes and legend settings to preserve chart formatting
        def _snap_axis(ax):
            if ax is None:
                return None
            d = {}
            try:
                d.update(
                    has_title=getattr(ax, "has_title", None),
                    title_text=(
                        ax.axis_title.text_frame.text
                        if getattr(ax, "has_title", False)
                        else None
                    ),
                    major_tick_mark=getattr(ax, "major_tick_mark", None),
                    minor_tick_mark=getattr(ax, "minor_tick_mark", None),
                    minimum_scale=getattr(ax, "minimum_scale", None),
                    maximum_scale=getattr(ax, "maximum_scale", None),
                    major_unit=getattr(ax, "major_unit", None),
                    minor_unit=getattr(ax, "minor_unit", None),
                    reverse_order=getattr(ax, "reverse_order", None),
                    tick_label_position=getattr(ax, "tick_label_position", None),
                    has_major_gridlines=getattr(ax, "has_major_gridlines", None),
                    has_minor_gridlines=getattr(ax, "has_minor_gridlines", None),
                    tick_num_fmt=(
                        getattr(getattr(ax, "tick_labels", None), "number_format", None)
                        if hasattr(ax, "tick_labels")
                        else None
                    ),
                    tick_num_fmt_linked=(
                        getattr(
                            getattr(ax, "tick_labels", None),
                            "number_format_is_linked",
                            None,
                        )
                        if hasattr(ax, "tick_labels")
                        else None
                    ),
                    visible=getattr(ax, "visible", None),
                )
            except Exception:
                pass
            return d

        def _restore_axis(ax, d):
            if ax is None or not d:
                return
            try:
                if d.get("has_title") is not None:
                    ax.has_title = d["has_title"]
                    if d["has_title"] and d.get("title_text") is not None:
                        try:
                            ax.axis_title.text_frame.text = d["title_text"]
                        except Exception:
                            pass
                for k in (
                    "major_tick_mark",
                    "minor_tick_mark",
                    "minimum_scale",
                    "maximum_scale",
                    "major_unit",
                    "minor_unit",
                    "reverse_order",
                    "tick_label_position",
                    "has_major_gridlines",
                    "has_minor_gridlines",
                    "visible",
                ):
                    if d.get(k) is not None:
                        try:
                            setattr(ax, k, d[k])
                        except Exception:
                            pass
                tl = getattr(ax, "tick_labels", None)
                if tl is not None:
                    if d.get("tick_num_fmt") is not None:
                        tl.number_format = d["tick_num_fmt"]
                    if d.get("tick_num_fmt_linked") is not None:
                        tl.number_format_is_linked = d["tick_num_fmt_linked"]
            except Exception:
                pass

        def _snapshot_axes(ch):
            snap = {}
            for name in ("category_axis", "value_axis"):
                try:
                    ax = getattr(ch, name)
                except Exception:
                    ax = None
                snap[name] = _snap_axis(ax)
            return snap

        def _restore_axes(ch, snap):
            for name in ("category_axis", "value_axis"):
                try:
                    ax = getattr(ch, name)
                except Exception:
                    ax = None
                _restore_axis(ax, snap.get(name))

        # legend snapshot
        def _snapshot_legend(ch):
            try:
                if not getattr(ch, "has_legend", False):
                    return dict(has_legend=False)
                lg = ch.legend
                return dict(
                    has_legend=True,
                    position=getattr(lg, "position", None),
                    include_in_layout=getattr(lg, "include_in_layout", None),
                )
            except Exception:
                return {}

        def _restore_legend(ch, d):
            try:
                if d.get("has_legend") is not None:
                    ch.has_legend = d["has_legend"]
                if d.get("has_legend"):
                    lg = ch.legend
                    if d.get("position") is not None:
                        lg.position = d["position"]
                    if d.get("include_in_layout") is not None:
                        lg.include_in_layout = d["include_in_layout"]
            except Exception:
                pass

        axes_snapshot = _snapshot_axes(chart)
        legend_snapshot = _snapshot_legend(chart)

        cd = None
        if self.category_data is not None:
            cd = self.category_data.to_chart_data()
        elif self.xy_data is not None:
            cd = self.xy_data.to_chart_data()
        elif self.bubble_data is not None:
            cd = self.bubble_data.to_chart_data()
        if cd is not None:
            try:
                chart.replace_data(cd)
            except Exception:
                pass
            else:
                _restore_dlbls(chart, dl_snapshot)
                _restore_axes(chart, axes_snapshot)
                _restore_legend(chart, legend_snapshot)
        if self.style is not None:
            try:
                chart.chart_style = int(self.style)
            except Exception:
                pass
        if self.has_legend is not None:
            try:
                chart.has_legend = bool(self.has_legend)
            except Exception:
                pass
        if self.title is not None:
            try:
                chart.has_title = True
                chart.chart_title.text_frame.text = self.title
            except Exception:
                pass
