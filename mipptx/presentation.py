from __future__ import annotations

from io import BytesIO
from typing import List, Optional

from pptx import Presentation as _Presentation
from pydantic import ConfigDict, Field

from .base import JsonModel
from .charts import ChartModel
from .slide import SlideModel
from .utils import emu_to_inches, inches_to_emu


class PresentationModel(JsonModel):
    model_config = ConfigDict(extra="forbid")

    slide_width_in: Optional[float] = Field(
        None, gt=0, description="Slide width in inches"
    )
    slide_height_in: Optional[float] = Field(
        None, gt=0, description="Slide height in inches"
    )
    slides: List[SlideModel] = Field(default_factory=list)

    # ---------- Construction from python-pptx ----------
    @classmethod
    def from_presentation(cls, prs) -> "PresentationModel":
        """Build a model from a `pptx.Presentation` instance."""
        sw = emu_to_inches(getattr(prs, "slide_width", None))
        sh = emu_to_inches(getattr(prs, "slide_height", None))
        slides = [SlideModel.from_pptx(s) for s in prs.slides]
        return cls(slide_width_in=sw, slide_height_in=sh, slides=slides)

    # ---------- Apply/Build with python-pptx ----------
    def build_presentation(self, template: str | None = None):
        """Create a new `pptx.Presentation` from this model.

        If `template` is provided, it is used as the base presentation; otherwise
        python-pptx's default template is used.
        """
        prs = _Presentation(template) if template else _Presentation()
        # size
        if self.slide_width_in is not None:
            prs.slide_width = inches_to_emu(self.slide_width_in)
        if self.slide_height_in is not None:
            prs.slide_height = inches_to_emu(self.slide_height_in)

        # add slides
        for slide_model in self.slides:
            idx = getattr(slide_model, "layout_index", None) or 0
            try:
                layout = prs.slide_layouts[idx]
            except Exception:
                layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            slide_model.apply_to_pptx(slide)
        return prs

    def save(self, path: str, template: str | None = None) -> None:
        prs = self.build_presentation(template=template)
        prs.save(path)

    # Update only chart data and simple chart properties on an existing Presentation
    # This preserves layouts, themes, and non-chart content. Slides count and order must match or
    # extras are ignored. Returns the same `prs` object for chaining.
    def update_charts_in_presentation(self, prs, *, strict: bool = False):
        # iterate slides in parallel up to min length
        max_idx = min(len(self.slides), len(prs.slides))
        for i in range(max_idx):
            model_slide = self.slides[i]
            pptx_slide = prs.slides[i]

            # collect chart models in model order
            model_charts = [
                sh for sh in model_slide.shapes if isinstance(sh, ChartModel)
            ]
            # collect actual chart shapes in slide order
            slide_chart_shapes = [
                sh for sh in pptx_slide.shapes if getattr(sh, "has_chart", False)
            ]

            # try name-based pairing first, fallback to order-based
            by_name = {
                getattr(g, "name", None): g
                for g in slide_chart_shapes
                if getattr(g, "name", None)
            }
            used = set()
            pairs: list[tuple[ChartModel, object]] = []
            for cm in model_charts:
                g = by_name.get(getattr(cm, "name", None))
                if g is not None:
                    pairs.append((cm, g))
                    used.add(id(g))
                else:
                    pairs.append((cm, None))
            # fill Nones by order for remaining shapes
            order_shapes = [g for g in slide_chart_shapes if id(g) not in used]
            oi = 0
            filled: list[tuple[ChartModel, object]] = []
            for cm, g in pairs:
                if g is None:
                    if oi < len(order_shapes):
                        g = order_shapes[oi]
                        oi += 1
                    else:
                        continue
                filled.append((cm, g))

            for cm, gframe in filled:
                try:
                    if strict:
                        ok = cm.update_workbook_only(gframe.chart)
                        if not ok:
                            cm.apply_to_existing_chart(gframe.chart)
                    else:
                        cm.apply_to_existing_chart(gframe.chart)
                except Exception:
                    pass

        return prs
